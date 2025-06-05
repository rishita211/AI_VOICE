 from pptx import Presentation

def extract_text(pptx_file):



Extracts and returns text from a PPTX file.

prs Presentation(pptx_file)

full text = ""

for slide in prs.slides:

for shape in slide.shapes:

if hasattr(shape, "text"):

full_text + shape.text+

return full text.strip()

# Optional test run

if_name="main":

text extract_text("inputs/sample.pptx")
 print("Extracted text:\n", text)